import os
import pandas as pd
import numpy as np
from scipy.stats import norm

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


ANNOTATIONS = [
    {'year': '2011', 'gdp': '6,17%', 'type': 'green', 'date': '05/02/2012', 'x_pct': 0.08, 'y_pos': 1.25},
    {'year': '2012', 'gdp': '6,03%', 'type': 'green', 'date': '05/02/2013', 'x_pct': 0.14, 'y_pos': 1.35},
    {'year': '2013', 'gdp': '5,56%', 'type': 'green', 'date': '05/02/2014', 'x_pct': 0.20, 'y_pos': 1.55},
    {'year': '2014', 'gdp': '5,01%', 'type': 'green', 'date': '05/02/2015', 'x_pct': 0.26, 'y_pos': 1.75},
    {'year': '2015', 'gdp': '4,88%', 'type': 'red',   'date': '05/02/2016', 'x_pct': 0.32, 'y_pos': 2.30},
    {'year': '2016', 'gdp': '5,03%', 'type': 'green', 'date': '05/02/2017', 'x_pct': 0.38, 'y_pos': 1.75},
    {'year': '2017', 'gdp': '5,07%', 'type': 'green', 'date': '05/02/2018', 'x_pct': 0.44, 'y_pos': 1.75},
    {'year': '2018', 'gdp': '5,17%', 'type': 'green', 'date': '05/02/2019', 'x_pct': 0.50, 'y_pos': 1.75},
    {'year': '2019', 'gdp': '5,02%', 'type': 'green', 'date': '05/02/2020', 'x_pct': 0.56, 'y_pos': 1.75},
    {'year': '2020', 'gdp': '-2,07%', 'type': 'red',  'date': '05/02/2021', 'x_pct': 0.62, 'y_pos': 5.80},
    {'year': '2021', 'gdp': '3,07%', 'type': 'red',   'date': '05/02/2022', 'x_pct': 0.68, 'y_pos': 2.90},
    {'year': '2022', 'gdp': '5,31%', 'type': 'green', 'date': '05/02/2023', 'x_pct': 0.74, 'y_pos': 1.65},
    {'year': '2023', 'gdp': '5,05%', 'type': 'green', 'date': '05/02/2024', 'x_pct': 0.80, 'y_pos': 1.75},
    {'year': '2024', 'gdp': '5,03%', 'type': 'green', 'date': '05/02/2025', 'x_pct': 0.86, 'y_pos': 1.75},
    {'year': '2025', 'gdp': '5,11%', 'type': 'green', 'date': '05/02/2026', 'x_pct': 0.91, 'y_pos': 1.70},
    {'year': '2026*', 'gdp': '5,61%', 'type': 'green','date': '05/05/2026', 'x_pct': 0.96, 'y_pos': 1.45},
]


def build_actual_release_series(dates_list, df_prediction=None, horizon='Tahunan'):
    """
    Populates GDP release values ONLY on official release dates:
    - Yearly ('Tahunan'): 5 February (or next weekday if weekend/holiday).
    - Quarterly ('Triwulanan'): 5 February, 5 May, 5 August, 5 November (or next weekday if weekend/holiday).
    For all other daily prediction dates, returns None (empty cell in Excel Column C).
    """
    dates_dt = pd.to_datetime(dates_list, format='%d/%m/%Y', errors='coerce')
    if dates_dt.isna().all():
        dates_dt = pd.to_datetime(dates_list, errors='coerce')
        
    actual_dots = [None] * len(dates_list)
    if dates_dt.dropna().empty:
        return actual_dots
        
    min_yr = dates_dt.dropna().min().year
    max_yr = dates_dt.dropna().max().year
    
    release_months = [2] if 'tahunan' in horizon.lower() else [2, 5, 8, 11]
    
    target_release_dts = []
    for y in range(min_yr, max_yr + 1):
        for m in release_months:
            dt = pd.Timestamp(year=y, month=m, day=5)
            if dt.dayofweek == 5: # Saturday -> Monday 7th
                dt += pd.Timedelta(days=2)
            elif dt.dayofweek == 6: # Sunday -> Monday 6th
                dt += pd.Timedelta(days=1)
            target_release_dts.append(dt)
            
    rel_ann_map = {ann['date']: float(ann['gdp'].replace(',', '.').replace('%', '')) for ann in ANNOTATIONS}
    
    for trg_dt in target_release_dts:
        valid_indices = [idx for idx, d in enumerate(dates_dt) if pd.notna(d) and d >= trg_dt]
        if valid_indices:
            idx = valid_indices[0]
            dt_str = dates_list[idx]
            if df_prediction is not None and 'Actual' in df_prediction.columns:
                val = df_prediction.iloc[idx]['Actual']
                if pd.notna(val):
                    actual_dots[idx] = round(float(val), 2)
            elif dt_str in rel_ann_map:
                actual_dots[idx] = rel_ann_map[dt_str]
                
    return actual_dots


def add_native_dual_axis_slide(prs, slide_title, dates_list, prob_values, line_color_rgb, bins=(5.0, 5.35), slide_num=1, df_prediction=None, horizon='Tahunan'):
    """
    Adds a single slide with a NATIVE EDITABLE PowerPoint Combo Chart object:
    - Left Axis (LHS): Probability Line Plot (0-100%)
    - Right Axis (RHS): Actual Realized GDP Growth Dots & 5.0% / 5.35% Target Border Lines
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_title
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Calibri'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    
    b0, b1 = bins[0], bins[1]
    
    actual_dots = build_actual_release_series(dates_list, df_prediction=df_prediction, horizon=horizon)
    border_b0 = [b0] * len(dates_list)
    border_b1 = [b1] * len(dates_list)
    
    chart_data = CategoryChartData()
    chart_data.categories = dates_list
    chart_data.add_series('Probabilitas (LHS, %)', prob_values)
    chart_data.add_series('Realisasi GDP Dot (RHS, %)', actual_dots)
    chart_data.add_series(f'Batas Target {b0}% (RHS)', border_b0)
    chart_data.add_series(f'Batas Target {b1}% (RHS)', border_b1)
    
    # Native Line Chart
    x, y, cx, cy = Inches(0.5), Inches(1.1), Inches(12.333), Inches(5.8)
    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data)
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    
    # Disable gridlines on Primary axes
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False

    # Separate series 1, 2, 3 into secondary line chart referencing RHS vertical axis (axPos r)
    plot_area = chart._chartSpace.chart.plotArea
    line_chart = plot_area.xpath('c:lineChart')[0]
    
    cat_ax_id = line_chart.xpath('c:axId')[0].get('val')
    val_ax_sec_id = '50000002'
    
    ser_nodes = line_chart.xpath('c:ser')
    line_chart_sec = parse_xml(f'''
        <c:lineChart {nsdecls("c")}>
            <c:grouping val="standard"/>
            <c:axId val="{cat_ax_id}"/>
            <c:axId val="{val_ax_sec_id}"/>
        </c:lineChart>
    ''')
    
    for ser in ser_nodes[1:]:
        line_chart_sec.append(ser)
        
    line_chart.addnext(line_chart_sec)
    
    # Secondary value axis without gridlines
    val_ax_sec = parse_xml(f'''
        <c:valAx {nsdecls("c", "a")}>
            <c:axId val="{val_ax_sec_id}"/>
            <c:scaling>
                <c:orientation val="minMax"/>
            </c:scaling>
            <c:delete val="0"/>
            <c:axPos val="r"/>
            <c:crossAx val="{cat_ax_id}"/>
            <c:crosses val="max"/>
        </c:valAx>
    ''')
    
    plot_area.append(val_ax_sec)
    
    # Series formatting
    s0 = chart.series[0]
    s0.format.line.color.rgb = line_color_rgb
    s0.format.line.width = Pt(2.0)
    
    s1 = chart.series[1]
    s1.format.line.fill.background()
    
    s2 = chart.series[2]
    s2.format.line.color.rgb = RGBColor(183, 28, 28)
    s2.format.line.width = Pt(1.5)
    
    s3 = chart.series[3]
    s3.format.line.color.rgb = RGBColor(21, 128, 61)
    s3.format.line.width = Pt(1.5)
    
    # Milestone Callout Boxes
    for ann in ANNOTATIONS:
        box_x = Inches(0.5 + ann['x_pct'] * 11.5)
        box_y = Inches(ann['y_pos'])
        box_w = Inches(0.55)
        box_h = Inches(0.45)
        
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_x, box_y, box_w, box_h)
        rect.fill.solid()
        if ann['type'] == 'green':
            rect.fill.fore_color.rgb = RGBColor(219, 237, 198)
            font_clr = RGBColor(27, 94, 32)
        else:
            rect.fill.fore_color.rgb = RGBColor(255, 224, 178)
            font_clr = RGBColor(183, 28, 28)
            
        rect.line.fill.background()
        
        tf_b = rect.text_frame
        tf_b.word_wrap = True
        p_b1 = tf_b.paragraphs[0]
        p_b1.text = ann['year']
        p_b1.alignment = PP_ALIGN.CENTER
        p_b1.font.name = 'Calibri'
        p_b1.font.size = Pt(8.5)
        p_b1.font.bold = True
        p_b1.font.color.rgb = font_clr
        
        p_b2 = tf_b.add_paragraph()
        p_b2.text = ann['gdp']
        p_b2.alignment = PP_ALIGN.CENTER
        p_b2.font.name = 'Calibri'
        p_b2.font.size = Pt(8.5)
        p_b2.font.bold = True
        p_b2.font.color.rgb = font_clr

    # Footer Bar & Slide Number
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0, 32, 96)
    footer.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.4))
    p_num = num_box.text_frame.paragraphs[0]
    p_num.text = str(slide_num)
    p_num.alignment = PP_ALIGN.RIGHT
    p_num.font.name = 'Calibri'
    p_num.font.size = Pt(14)
    p_num.font.bold = True
    p_num.font.color.rgb = RGBColor(255, 255, 255)


def build_3slide_deck(df_data, horizon_label, output_pptx_filename, bins=(5.0, 5.35), df_prediction=None):
    """
    Builds a 3-slide PowerPoint presentation deck with NATIVE EDITABLE Dual-Axis Combo Chart objects.
    """
    if not PPTX_AVAILABLE:
        print("Warning: python-pptx is not installed. Skipping PPTX generation.")
        return False
        
    b0, b1 = bins[0], bins[1]
    dates = df_data['Day Prediction'].tolist()
    
    col_b0 = f'P(< {b0}%)'
    col_b1 = f'P({b0}% - {b1}%)'
    col_b2 = f'P(> {b1}%)'
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Below b0
    t1 = f"Probabilitas Pertumbuhan Ekonomi {horizon_label} di Bawah {b0}%"
    add_native_dual_axis_slide(prs, t1, dates, df_data[col_b0].tolist(), RGBColor(230, 81, 0), bins=bins, slide_num=1, df_prediction=df_prediction, horizon=horizon_label)
    
    # Slide 2: Target Band b0 - b1
    t2 = f"Probabilitas Pertumbuhan Ekonomi {horizon_label} {b0}% - {b1}%"
    add_native_dual_axis_slide(prs, t2, dates, df_data[col_b1].tolist(), RGBColor(245, 158, 11), bins=bins, slide_num=2, df_prediction=df_prediction, horizon=horizon_label)
    
    # Slide 3: Above b1
    t3 = f"Probabilitas Pertumbuhan Ekonomi {horizon_label} di Atas {b1}%"
    add_native_dual_axis_slide(prs, t3, dates, df_data[col_b2].tolist(), RGBColor(16, 185, 129), bins=bins, slide_num=3, df_prediction=df_prediction, horizon=horizon_label)
    
    os.makedirs(os.path.dirname(output_pptx_filename), exist_ok=True)
    prs.save(output_pptx_filename)
    print(f"Successfully generated 3-slide NATIVE EDITABLE Dual-Axis PPTX deck: {output_pptx_filename}")
    return True


def generate_probability_pptx(df_prediction, std_annual=None, std_nowcast=None, bins=(5.0, 5.35), output_dir="output"):
    """
    Generates 3-slide PowerPoint presentation decks (.pptx) with NATIVE EDITABLE Dual-Axis Combo Charts for FY and Quarterly growth probabilities.
    """
    if not PPTX_AVAILABLE:
        return
        
    b0, b1 = bins[0], bins[1]
    dates = pd.to_datetime(df_prediction['Day Prediction']).dt.strftime('%d/%m/%Y').tolist()
    ref_q = pd.to_datetime(df_prediction['Reference Quarter']).dt.strftime('%d/%m/%Y').tolist()
    
    # 1. Full Year (Annual) 3-Slide Native Editable Deck
    if 'Annual Nowcast' in df_prediction.columns and std_annual is not None:
        mu_ann = df_prediction['Annual Nowcast'].values
        std_ann = np.squeeze(np.asarray(std_annual))
        if len(std_ann) < len(mu_ann):
            std_ann = np.pad(std_ann, (0, len(mu_ann) - len(std_ann)), mode='edge')
        elif len(std_ann) > len(mu_ann):
            std_ann = std_ann[:len(mu_ann)]
            
        p_ann_b0 = norm.cdf(b0, loc=mu_ann, scale=std_ann) * 100.0
        p_ann_b1 = (norm.cdf(b1, loc=mu_ann, scale=std_ann) - norm.cdf(b0, loc=mu_ann, scale=std_ann)) * 100.0
        p_ann_b2 = (1.0 - norm.cdf(b1, loc=mu_ann, scale=std_ann)) * 100.0
        
        df_fy = pd.DataFrame({
            'Day Prediction': dates,
            'Reference Quarter': ref_q,
            f'P(< {b0}%)': p_ann_b0.round(2),
            f'P({b0}% - {b1}%)': p_ann_b1.round(2),
            f'P(> {b1}%)': p_ann_b2.round(2)
        })
        
        pptx_fy = os.path.join(output_dir, 'Probabilitas_Pertumbuhan_Ekonomi_Tahunan_3Slides_Native.pptx')
        build_3slide_deck(df_fy, 'Tahunan', pptx_fy, bins=bins, df_prediction=df_prediction)
        
    # 2. Quarterly 3-Slide Native Editable Deck
    if 'Nowcast' in df_prediction.columns and std_nowcast is not None:
        mu_now = df_prediction['Nowcast'].values
        std_now = np.squeeze(np.asarray(std_nowcast))
        if len(std_now) < len(mu_now):
            std_now = np.pad(std_now, (0, len(mu_now) - len(std_now)), mode='edge')
        elif len(std_now) > len(mu_now):
            std_now = std_now[:len(mu_now)]
            
        p_now_b0 = norm.cdf(b0, loc=mu_now, scale=std_now) * 100.0
        p_now_b1 = (norm.cdf(b1, loc=mu_now, scale=std_now) - norm.cdf(b0, loc=mu_now, scale=std_now)) * 100.0
        p_now_b2 = (1.0 - norm.cdf(b1, loc=mu_now, scale=std_now)) * 100.0
        
        df_q = pd.DataFrame({
            'Day Prediction': dates,
            'Reference Quarter': ref_q,
            f'P(< {b0}%)': p_now_b0.round(2),
            f'P({b0}% - {b1}%)': p_now_b1.round(2),
            f'P(> {b1}%)': p_now_b2.round(2)
        })
        
        pptx_q = os.path.join(output_dir, 'Probabilitas_Pertumbuhan_Ekonomi_Triwulanan_3Slides_Native.pptx')
        build_3slide_deck(df_q, 'Triwulanan', pptx_q, bins=bins, df_prediction=df_prediction)
